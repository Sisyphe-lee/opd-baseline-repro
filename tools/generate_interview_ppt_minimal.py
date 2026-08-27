#!/usr/bin/env python3
"""Generate a six-slide, fully editable minimalist interview deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from generate_interview_ppt_editable import (
    FONT,
    FONT_LATIN,
    FONT_MATH,
    INK,
    INK_2,
    LINE,
    MUTED,
    PANEL,
    PLUM,
    PLUM_LIGHT,
    REFERENCE,
    ROOT,
    WHITE,
    connector,
    ellipse,
    fill,
    line_style,
    rect,
    remove_all_slides,
    rich,
    text,
)


OUT_DIR = ROOT / "deliverables" / "interview_self_intro_opd"
TOTAL = 6


def header(slide, title: str, number: int, subtitle: str | None = None) -> None:
    text(slide, (64, 30, 1536, 82), title, 27, INK, True)
    if subtitle:
        text(slide, (64, 83, 1536, 112), subtitle, 12.5, MUTED)
    connector(slide, 64, 126, 1536, 126, LINE, 1.0)
    text(
        slide,
        (1450, 844, 1536, 870),
        f"{number:02d} / {TOTAL:02d}",
        10,
        MUTED,
        False,
        PP_ALIGN.RIGHT,
        font_name=FONT_LATIN,
    )


def label(slide, x, y, value, color=PLUM) -> None:
    text(slide, (x, y, x + 220, y + 28), value, 11, color, True, font_name=FONT_LATIN)


def divider(slide, x, y1, y2) -> None:
    connector(slide, x, y1, x, y2, LINE, 1.0)


def bullet(slide, x, y, width, value, size=15, color=INK_2) -> None:
    ellipse(slide, (x, y + 8, x + 8, y + 16), PLUM, None)
    text(slide, (x + 22, y, x + width, y + 40), value, size, color)


def link_text(slide, box, value, url, color=PLUM) -> None:
    text(slide, box, value, 11.5, color, True, url=url)


def arrow(slide, x1, y, x2) -> None:
    connector(slide, x1, y, x2 - 12, y, PLUM_LIGHT, 1.8)
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches((x2 - 18) / 120),
        Inches((y - 8) / 120),
        Inches(16 / 120),
        Inches(16 / 120),
    )
    tri.rotation = 90
    fill(tri, PLUM_LIGHT)
    line_style(tri, None)


def slide01(slide) -> None:
    # One-page personal introduction; no separate cover or project-overview page.
    text(slide, (64, 42, 980, 102), "孙家兴", 34, INK, True)
    text(slide, (64, 108, 980, 142), "北京大学智能学院 · 硕士研究生", 16, MUTED)
    connector(slide, 64, 170, 1536, 170, PLUM, 2.0)

    label(slide, 64, 215, "PROFILE")
    text(slide, (64, 254, 690, 294), "教育与研究方向", 19, INK, True)
    bullet(slide, 64, 318, 620, "北京大学智能学院，硕士（2025.09—至今）", 14)
    bullet(slide, 64, 365, 620, "北京大学信息科学技术学院，本科", 14)
    bullet(slide, 64, 412, 620, "大模型后训练、Agentic Learning、采样与决策", 14)

    divider(slide, 760, 220, 480)
    label(slide, 820, 215, "SKILLS")
    text(slide, (820, 254, 1470, 294), "技术能力", 19, INK, True)
    text(slide, (820, 318, 1480, 410), "Python / PyTorch / veRL / vLLM\nFSDP / Ray / ALFWorld / 自动评测", 15, INK_2)
    text(slide, (820, 426, 1480, 465), "研究方法：可复现基线 · 机制诊断 · 配对评测", 13, MUTED)

    connector(slide, 64, 515, 1536, 515, LINE, 1.0)
    label(slide, 64, 548, "CORE RESEARCH")
    text(slide, (64, 587, 700, 624), "01  熵自适应策略蒸馏", 19, INK, True)
    text(slide, (64, 635, 700, 677), "为每条轨迹动态选择可靠蒸馏前缀", 14, MUTED)
    text(slide, (64, 690, 260, 737), "86.86%", 27, PLUM, True, font_name=FONT_LATIN)
    text(slide, (270, 704, 580, 734), "ALFWorld full274", 11, MUTED, font_name=FONT_LATIN)
    link_text(slide, (64, 755, 250, 782), "GitHub ↗", "https://github.com/Sisyphe-lee/opd-baseline-repro")

    divider(slide, 760, 555, 790)
    text(slide, (820, 587, 1480, 624), "02  SoftSat：采样失效修复", 19, INK, True)
    text(slide, (820, 635, 1480, 677), "解释 Power Sampling 失效并校准锐化强度", 14, MUTED)
    text(slide, (820, 690, 1000, 737), "7.42×", 27, PLUM, True, font_name=FONT_LATIN)
    text(slide, (1010, 704, 1320, 734), "加权推理加速", 11, MUTED)
    link_text(slide, (820, 755, 1005, 782), "GitHub ↗", "https://github.com/jiaxingsunpku/SoftSat")

    text(slide, (64, 832, 1300, 860), "2501213407@stu.pku.edu.cn   ·   +86 188-1171-7335", 10.5, MUTED, font_name=FONT_LATIN)
    text(slide, (1450, 844, 1536, 870), "01 / 06", 10, MUTED, False, PP_ALIGN.RIGHT, font_name=FONT_LATIN)


def slide02(slide) -> None:
    header(slide, "研究问题：长程交互使蒸馏监督逐步失真", 2)
    text(slide, (64, 160, 1536, 204), "早期动作不仅影响后续文本，还会直接改变环境状态。", 19, INK, True)

    nodes = [
        ("Student action", "动作偏差"),
        ("State transition", "状态改变"),
        ("Observation", "上下文漂移"),
        ("Teacher", "不确定性上升"),
    ]
    xs = [76, 450, 824, 1198]
    for i, (x, (en, zh)) in enumerate(zip(xs, nodes)):
        ellipse(slide, (x, 290, x + 54, 344), PLUM if i == 0 else PANEL, PLUM if i else None)
        text(slide, (x + 72, 284, x + 300, 316), en, 13, INK, True, font_name=FONT_LATIN)
        text(slide, (x + 72, 327, x + 300, 358), zh, 13, MUTED)
        if i < 3:
            arrow(slide, x + 300, 320, xs[i + 1] - 18)

    connector(slide, 64, 430, 1536, 430, LINE, 1.0)
    label(slide, 64, 472, "WHY IT MATTERS")
    bullet(slide, 64, 520, 1430, "后半段的低质量 dense supervision 可能从帮助变成噪声。", 16)
    bullet(slide, 64, 578, 1430, "固定截断会同时错杀仍可学习的轨迹，并放过已经漂移的轨迹。", 16)

    rect(slide, (64, 670, 1536, 765), PANEL, None, False)
    text(slide, (92, 694, 280, 730), "核心问题", 13, PLUM, True)
    text(slide, (280, 686, 1490, 742), "能否为每条 trajectory 单独判断 Teacher 监督何时不再可靠？", 20, INK, True, valign=MSO_ANCHOR.MIDDLE)


def slide03(slide) -> None:
    header(slide, "已有方法：TCOD 提升基线，但仍使用统一训练深度", 3)
    # Minimal three-column comparison, avoiding decorative cards.
    cols = [64, 515, 966, 1536]
    for x in cols[1:-1]:
        divider(slide, x, 195, 660)
    label(slide, 64, 182, "VANILLA OPD", MUTED)
    text(slide, (64, 230, 470, 272), "整条轨迹蒸馏", 20, INK, True)
    text(slide, (64, 302, 470, 346), "所有 turn 都进入 loss", 14, MUTED)
    text(slide, (64, 430, 470, 492), "79.56%", 35, PLUM_LIGHT, True, font_name=FONT_LATIN)
    text(slide, (64, 503, 470, 536), "218 / 274", 12, MUTED, font_name=FONT_LATIN)

    label(slide, 565, 182, "TCOD-F2B", MUTED)
    text(slide, (565, 230, 920, 272), "逐步扩展 horizon", 20, INK, True)
    text(slide, (565, 302, 920, 366), "训练越深入，统一增加所有轨迹的 K", 14, MUTED)
    text(slide, (565, 430, 920, 492), "84.67%", 35, PLUM, True, font_name=FONT_LATIN)
    text(slide, (565, 503, 920, 536), "232 / 274", 12, MUTED, font_name=FONT_LATIN)

    label(slide, 1016, 182, "LIMITATION", PLUM)
    text(slide, (1016, 230, 1500, 272), "全局 schedule", 20, INK, True)
    bullet(slide, 1016, 302, 475, "同一阶段，所有轨迹使用相同 K", 14)
    bullet(slide, 1016, 358, 475, "无法感知每条轨迹的实际漂移时刻", 14)
    bullet(slide, 1016, 414, 475, "训练进度不能替代轨迹状态", 14)

    rect(slide, (64, 700, 1536, 770), PLUM, None, False)
    text(slide, (92, 715, 1500, 754), "改进方向：从 global curriculum 转向 trajectory-specific loss selection。", 18, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
    link_text(slide, (64, 806, 520, 832), "TCOD · arXiv:2604.24005 ↗", "https://arxiv.org/abs/2604.24005", MUTED)


def slide04(slide) -> None:
    header(slide, "我们的想法：为每条轨迹定位可蒸馏边界", 4)
    text(slide, (64, 160, 1536, 200), "使用 Teacher entropy 的相对漂移，而不是固定 turn 数。", 18, INK, True)

    # Editable entropy curve.
    x0, y0, x1, y1 = 90, 650, 850, 275
    connector(slide, x0, y0, x1, y0, MUTED, 1.0)
    connector(slide, x0, y0, x0, y1, MUTED, 1.0)
    vals = [0.18, 0.20, 0.19, 0.23, 0.26, 0.32, 0.42, 0.55, 0.63]
    pts = []
    for i, value in enumerate(vals):
        pts.append((x0 + 40 + i * 84, y0 - 35 - value * 470))
    for a, b in zip(pts, pts[1:]):
        connector(slide, *a, *b, PLUM, 2.6)
    for x, y in pts:
        ellipse(slide, (x - 5, y - 5, x + 5, y + 5), PLUM, None)
    frontier_x = pts[6][0]
    connector(slide, frontier_x, y1 + 5, frontier_x, y0, PLUM_LIGHT, 1.8)
    text(slide, (frontier_x + 12, y1 + 10, frontier_x + 180, y1 + 40), "frontier  fᵢ", 12, PLUM, True, font_name=FONT_MATH)
    text(slide, (135, 670, frontier_x - 10, 702), "保留 loss", 12, PLUM, True, PP_ALIGN.CENTER)
    text(slide, (frontier_x + 10, 670, x1, 702), "屏蔽 suffix", 12, MUTED, False, PP_ALIGN.CENTER)

    divider(slide, 930, 235, 740)
    label(slide, 990, 240, "ALGORITHM")
    text(slide, (990, 290, 1480, 332), "1  建立局部基线", 17, INK, True)
    text(slide, (1020, 338, 1480, 378), "首 3 turn：Bᵢ = ⅓ΣHᵢ,ₜ", 14, MUTED, font_name=FONT_MATH)
    text(slide, (990, 416, 1480, 458), "2  检测持续漂移", 17, INK, True)
    text(slide, (1020, 464, 1480, 516), "连续 3 turn 平均漂移 ≥ τ", 14, MUTED)
    text(slide, (990, 554, 1480, 596), "3  选择训练前缀", 17, INK, True)
    text(slide, (1020, 602, 1480, 654), "仅保留 t < fᵢ；无触发则保留完整轨迹", 14, MUTED)

    connector(slide, 64, 770, 1536, 770, LINE, 1.0)
    text(slide, (64, 798, 1536, 830), "实验控制：完整 rollout 和 Teacher 评分不变，只改变进入 loss 的数据。", 13, INK_2, True)


def slide05(slide) -> None:
    header(slide, "实验结果：τ = 0.10 达到 86.86%", 5, "ALFWorld full274 · 250 steps · seed 42 · 同一评测协议")

    # Large, simple editable bars.
    x_axis, y_base, y_top = 110, 700, 260
    connector(slide, x_axis, y_base, 900, y_base, MUTED, 1.0)
    methods = ["Vanilla OPD", "TCOD-F2B", "Entropy Adaptive"]
    rates = [79.56, 84.67, 86.86]
    colors = [LINE, PLUM_LIGHT, PLUM]
    for i, (name, rate, color) in enumerate(zip(methods, rates, colors)):
        x = 180 + i * 245
        height = (rate - 70) / 20 * (y_base - y_top)
        rect(slide, (x, y_base - height, x + 125, y_base), color, None, False)
        text(slide, (x - 30, y_base - height - 50, x + 155, y_base - height - 10), f"{rate:.2f}%", 18, PLUM if i == 2 else INK_2, True, PP_ALIGN.CENTER, font_name=FONT_LATIN)
        text(slide, (x - 40, y_base + 20, x + 165, y_base + 50), name, 11, MUTED, i == 2, PP_ALIGN.CENTER, font_name=FONT_LATIN)

    divider(slide, 955, 205, 770)
    label(slide, 1015, 205, "RESULT")
    text(slide, (1015, 255, 1480, 318), "+7.30 pp", 31, PLUM, True, font_name=FONT_LATIN)
    text(slide, (1015, 320, 1480, 350), "vs Vanilla OPD", 12, MUTED, font_name=FONT_LATIN)
    text(slide, (1015, 405, 1480, 468), "+2.19 pp", 31, PLUM, True, font_name=FONT_LATIN)
    text(slide, (1015, 470, 1480, 500), "vs TCOD-F2B", 12, MUTED, font_name=FONT_LATIN)
    text(slide, (1015, 568, 1200, 600), "Seen", 12, MUTED, font_name=FONT_LATIN)
    text(slide, (1210, 555, 1480, 605), "87.86%", 22, INK, True, font_name=FONT_LATIN)
    text(slide, (1015, 630, 1200, 662), "Unseen", 12, MUTED, font_name=FONT_LATIN)
    text(slide, (1210, 617, 1480, 667), "85.82%", 22, INK, True, font_name=FONT_LATIN)

    connector(slide, 64, 795, 1536, 795, LINE, 1.0)
    rich(slide, (64, 816, 1450, 844), [
        {"text": "阈值扫描  ", "color": MUTED, "bold": True, "size": 10.5},
        {"text": "τ=.05  81.75%   ·   τ=.075  80.29%   ·   ", "color": MUTED, "size": 10.5, "font": FONT_LATIN},
        {"text": "τ=.10  86.86%", "color": PLUM, "bold": True, "size": 10.5, "font": FONT_LATIN},
        {"text": "   ·   τ=.125  82.48%", "color": MUTED, "size": 10.5, "font": FONT_LATIN},
    ], 10.5)


def slide06(slide) -> None:
    header(slide, "结论：自适应筛选在训练早期介入，后期自然退场", 6)

    # Two numbers instead of two dense line charts.
    label(slide, 64, 190, "TRAINING DYNAMICS")
    text(slide, (64, 245, 500, 310), "48.0%  →  9.0%", 31, PLUM, True, font_name=FONT_LATIN)
    text(slide, (64, 315, 500, 352), "frontier 触发率", 13, MUTED)
    divider(slide, 590, 220, 380)
    text(slide, (680, 245, 1120, 310), "20.0  →  28.2", 31, PLUM, True, font_name=FONT_LATIN)
    text(slide, (680, 315, 1120, 352), "平均可蒸馏 horizon", 13, MUTED)
    text(slide, (1210, 250, 1510, 348), "模型越成熟，\n算法越少干预。", 19, INK, True, PP_ALIGN.RIGHT)

    connector(slide, 64, 425, 1536, 425, LINE, 1.0)
    label(slide, 64, 468, "TAKEAWAYS")
    bullet(slide, 64, 520, 1430, "相对熵漂移可以作为 trajectory-level 的可靠筛选信号。", 17)
    bullet(slide, 64, 585, 1430, "关键不是“熵越低越好”，而是过滤强度与长程状态覆盖的平衡。", 17)
    bullet(slide, 64, 650, 1430, "最强证据来自 final checkpoint；后续需要多种子复测和连续加权。", 17)

    rect(slide, (64, 755, 1536, 817), PANEL, None, False)
    text(slide, (92, 765, 1490, 807), "证据边界：对 TCOD 的 +2.19 pp 在单种子配对检验下尚不显著（McNemar p=.362）。", 13, MUTED, valign=MSO_ANCHOR.MIDDLE)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(REFERENCE))
    remove_all_slides(prs)
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[0]
    for fn in (slide01, slide02, slide03, slide04, slide05, slide06):
        slide = prs.slides.add_slide(blank)
        fn(slide)
    prs.core_properties.title = "孙家兴｜面试自我介绍｜简约可编辑版"
    prs.core_properties.subject = "Entropy-adaptive on-policy distillation"
    prs.core_properties.author = "孙家兴"
    out = OUT_DIR / "孙家兴_面试自我介绍_OPD项目_简约可编辑.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
