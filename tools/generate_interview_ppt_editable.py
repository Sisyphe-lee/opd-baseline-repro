#!/usr/bin/env python3
"""Generate a fully editable interview PowerPoint deck.

Every visible item is a native PowerPoint text box, shape, connector, or line.
No pictures are embedded.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "deliverables" / "interview_self_intro_opd"
REFERENCE = Path("/vepfs-mlp2/mlp-public/252302025/sjx/.codex-sjx/attachments/3147fb86-839f-4e03-b4c3-0730d2b8a9cd/on-policy-distillation-agentic-opd-revised.pptx")
W, H = 1600, 900
INK="211F20"; INK_2="343033"; MUTED="716A6E"; PLUM="70485B"; PLUM_LIGHT="9B7184"
ROSE="8C3E49"; GREEN="315C4A"; PAPER="FFFFFF"; PANEL="F4F0F2"; PANEL_2="EEE3E8"
LINE="D8D0D4"; GREEN_PANEL="E5EEE9"; WHITE="FFFFFF"
FONT="Microsoft YaHei"; FONT_LATIN="Aptos"; FONT_MATH="Cambria Math"

def emu(v): return Inches(v/120.0)
def rgb(c): return RGBColor.from_string(c.lstrip("#"))

def remove_all_slides(prs):
    for sid in list(prs.slides._sldIdLst):
        prs.part.drop_rel(sid.rId); prs.slides._sldIdLst.remove(sid)

def fill(shape, color):
    if color is None: shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb=rgb(color)

def line_style(shape, color, width=1.0):
    if color is None: shape.line.fill.background()
    else:
        shape.line.color.rgb=rgb(color); shape.line.width=Pt(width)

def rect(slide, box, fill_color=PANEL, line_color=LINE, rounded=True, line_width=1.0):
    x1,y1,x2,y2=box
    kind=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s=slide.shapes.add_shape(kind,emu(x1),emu(y1),emu(x2-x1),emu(y2-y1))
    fill(s,fill_color); line_style(s,line_color,line_width)
    if rounded:
        try: s.adjustments[0]=0.12
        except Exception: pass
    return s

def ellipse(slide, box, fill_color=PLUM, line_color=None, line_width=1.0):
    x1,y1,x2,y2=box
    s=slide.shapes.add_shape(MSO_SHAPE.OVAL,emu(x1),emu(y1),emu(x2-x1),emu(y2-y1))
    fill(s,fill_color); line_style(s,line_color,line_width); return s

def connector(slide,x1,y1,x2,y2,color=LINE,width=1.5):
    s=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,emu(x1),emu(y1),emu(x2),emu(y2))
    line_style(s,color,width); return s

def text(slide,box,value,size=20,color=INK,bold=False,align=PP_ALIGN.LEFT,valign=MSO_ANCHOR.TOP,font_name=FONT,margin=0,url=None):
    x1,y1,x2,y2=box
    s=slide.shapes.add_textbox(emu(x1),emu(y1),emu(x2-x1),emu(y2-y1))
    tf=s.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=emu(margin)
    tf.vertical_anchor=valign
    for i,ln in enumerate(str(value).split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_before=Pt(0); p.space_after=Pt(0); p.line_spacing=1.0
        r=p.add_run(); r.text=ln; r.font.name=font_name; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(color)
        if url: r.hyperlink.address=url
    return s

def rich(slide,box,segments,size=20,align=PP_ALIGN.LEFT,valign=MSO_ANCHOR.TOP):
    x1,y1,x2,y2=box
    s=slide.shapes.add_textbox(emu(x1),emu(y1),emu(x2-x1),emu(y2-y1))
    tf=s.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=valign
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=align; p.space_before=Pt(0); p.space_after=Pt(0)
    for seg in segments:
        r=p.add_run(); r.text=seg['text']; r.font.name=seg.get('font',FONT); r.font.size=Pt(seg.get('size',size)); r.font.bold=seg.get('bold',False); r.font.color.rgb=rgb(seg.get('color',INK))
        if seg.get('url'): r.hyperlink.address=seg['url']
    return s

def pill(slide,box,label,fill_color,color=INK,size=15,bold=False,url=None):
    s=rect(slide,box,fill_color,None,True)
    text(slide,box,label,size,color,bold,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE,margin=2,url=url)
    if url: s.click_action.hyperlink.address=url
    return s

def base(slide,title,number,kicker=None):
    text(slide,(58,25,1542,82),title,27,PLUM,True)
    if kicker: text(slide,(58,78,1542,105),kicker,12,MUTED)
    connector(slide,58,114,1542,114,LINE,1.2)
    text(slide,(1450,842,1542,872),f"{number:02d} / 10",10,MUTED,False,PP_ALIGN.RIGHT,font_name=FONT_LATIN)

def footer(slide,value): text(slide,(58,835,1350,864),value,10,MUTED)

def card(slide,box,index,title,body,accent=PLUM,fill_color=PANEL):
    x1,y1,x2,y2=box; rect(slide,box,fill_color,LINE,True,0.8)
    text(slide,(x1+28,y1+18,x2-20,y1+45),index,11,accent,True,font_name=FONT_LATIN)
    text(slide,(x1+28,y1+55,x2-20,y1+100),title,18,INK,True)
    text(slide,(x1+28,y1+110,x2-22,y2-18),body,13,MUTED)

def arrow(slide,x1,y,x2,color=PLUM_LIGHT):
    connector(slide,x1,y,x2-12,y,color,2.2)
    s=slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,emu(x2-18),emu(y-9),emu(18),emu(18)); s.rotation=90; fill(s,color); line_style(s,None)

def bullet(slide,x,y,width,label,size=15,color=INK_2,bullet_color=PLUM_LIGHT,bold=False):
    ellipse(slide,(x,y+7,x+10,y+17),bullet_color,None)
    text(slide,(x+24,y,x+width,y+40),label,size,color,bold)

def trajectory(slide,x,y,label,active,active_color,box_count=6):
    text(slide,(x,y-5,x+112,y+35),label,12,MUTED,font_name=FONT_LATIN)
    px=x+130
    for i in range(box_count):
        rect(slide,(px,y,px+70,y+38),active_color if i<active else LINE,None,True)
        px+=80

def mini_line(slide,box,values,color,labels,title,suffix="%"):
    x1,y1,x2,y2=box; rect(slide,box,WHITE,LINE,True,0.8); text(slide,(x1+28,y1+22,x2-20,y1+55),title,15,INK,True)
    px1,py1,px2,py2=x1+54,y2-70,x2-34,y1+95
    lo,hi=min(values),max(values); span=max(hi-lo,1e-6); pts=[]
    for i,v in enumerate(values):
        x=px1+i*(px2-px1)/(len(values)-1); y=py1-(v-lo)/span*(py1-py2); pts.append((x,y))
    for a,b in zip(pts,pts[1:]): connector(slide,*a,*b,color,2.4)
    for i,((x,y),v,lab) in enumerate(zip(pts,values,labels)):
        ellipse(slide,(x-6,y-6,x+6,y+6),color,None)
        if i in (0,len(values)-1): text(slide,(x-48,y-35,x+48,y-10),f"{v:.1f}{suffix}",11,color,True,PP_ALIGN.CENTER,font_name=FONT_LATIN)
        text(slide,(x-45,py1+13,x+45,py1+38),lab,9,MUTED,False,PP_ALIGN.CENTER,font_name=FONT_LATIN)

def slide01(slide):
    rect(slide,(0,0,38,H),PLUM,None,False); rect(slide,(1220,0,W,H),PANEL,None,False)
    ellipse(slide,(1300,95,1515,310),PLUM,None); text(slide,(1300,95,1515,310),"SJX",36,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE,FONT_LATIN)
    text(slide,(98,170,1100,255),"孙家兴",47,PLUM,True); text(slide,(98,282,1100,345),"面试自我介绍",30,INK,True)
    rich(slide,(98,380,1110,420),[
        {'text':'大模型后训练','color':MUTED,'size':17},{'text':'  ·  ','color':PLUM_LIGHT,'size':17,'font':FONT_LATIN},
        {'text':'Agentic Learning','color':MUTED,'size':17,'font':FONT_LATIN},{'text':'  ·  ','color':PLUM_LIGHT,'size':17,'font':FONT_LATIN},
        {'text':'推理增强','color':MUTED,'size':17}],17)
    connector(slide,98,478,1050,478,LINE,1.6)
    text(slide,(98,520,900,560),"北京大学 · 智能学院",17,INK_2,True)
    text(slide,(98,580,900,615),"M.S. Candidate  |  August 2026",13,MUTED,font_name=FONT_LATIN)
    pill(slide,(1270,400,1550,464),"Research × Engineering",WHITE,PLUM,12,True)
    text(slide,(1270,515,1550,610),"从问题定义、可复现基线，\n到机制诊断与完整评测",15,INK_2)
    text(slide,(98,812,900,846),"2501213407@stu.pku.edu.cn  ·  +86 188-1171-7335",11,MUTED,font_name=FONT_LATIN)
    text(slide,(1450,842,1542,872),"01 / 10",10,MUTED,False,PP_ALIGN.RIGHT,font_name=FONT_LATIN)

def slide02(slide):
    base(slide,"基本信息：研究型算法工程与可验证实验",2)
    rect(slide,(58,155,505,800),PLUM,None,True); text(slide,(94,190,465,248),"孙家兴",32,WHITE,True)
    text(slide,(94,270,465,310),"北京大学 · 智能学院",17,WHITE,True); text(slide,(94,315,465,350),"硕士研究生（2025.09—至今）",13,PANEL_2)
    connector(slide,94,372,465,372,PLUM_LIGHT,1.2); text(slide,(94,408,465,438),"教育背景",12,LINE)
    text(slide,(94,450,465,520),"北京大学 · 信息科学技术学院\n本科（2021.09—2025.06）",15,WHITE)
    text(slide,(94,568,465,598),"联系方式",12,LINE); text(slide,(94,610,465,670),"188-1171-7335\n2501213407@stu.pku.edu.cn",12,WHITE,font_name=FONT_LATIN)
    pill(slide,(94,720,318,770),"GitHub  ↗",WHITE,PLUM,13,True,"https://github.com/Sisyphe-lee")
    card(slide,(550,155,1015,440),"01","研究方向","在线策略蒸馏 · 长程 Agent\n采样与多轨迹决策 · 自动评测",PLUM,PANEL)
    card(slide,(1050,155,1542,440),"02","训练与系统","PyTorch · veRL · vLLM · FSDP\nRay · Tool Calling · JSONL artifacts",GREEN,GREEN_PANEL)
    card(slide,(550,475,1015,800),"03","我的工作方式","先冻结可比较的基线，再做诊断；\n把每个提升拆成可证伪的假设。",ROSE,PANEL_2)
    card(slide,(1050,475,1542,800),"04","相关经历","大模型蒸馏、代码 Agent、\nPPO / MAPPO 多智能体控制。",PLUM,PANEL)

def slide03(slide):
    base(slide,"两项核心研究：从数据机制到决策机制",3)
    text(slide,(58,142,1542,180),"共同主线：不是只追求一个更高分，而是解释模型为什么失效、什么信号能修复。",15,INK_2)
    rect(slide,(58,220,775,742),PANEL,PLUM,True,1.8); pill(slide,(92,252,255,302),"PROJECT 01",PLUM,WHITE,11,True)
    text(slide,(92,340,720,392),"熵自适应策略蒸馏",24,PLUM,True); text(slide,(92,404,720,436),"Agentic OPD / ALFWorld",13,MUTED,font_name=FONT_LATIN)
    connector(slide,92,460,720,460,LINE,1.0)
    text(slide,(92,496,155,530),"问题",12,PLUM,True); text(slide,(172,496,720,536),"固定 horizon 无法判断每条轨迹何时开始漂移",14,INK_2)
    text(slide,(92,552,155,586),"想法",12,PLUM,True); text(slide,(172,552,720,592),"以教师熵的相对变化定位 distillable frontier",14,INK_2)
    text(slide,(92,618,285,672),"86.86%",29,ROSE,True,font_name=FONT_LATIN); text(slide,(300,638,500,668),"ALFWorld full274",11,MUTED,font_name=FONT_LATIN)
    pill(slide,(556,662,720,712),"GitHub  ↗",WHITE,PLUM,12,True,"https://github.com/Sisyphe-lee/opd-baseline-repro")
    rect(slide,(825,220,1542,742),GREEN_PANEL,GREEN,True,1.8); pill(slide,(859,252,1022,302),"PROJECT 02",GREEN,WHITE,11,True)
    text(slide,(859,340,1487,392),"SoftSat：采样失效修复",23,GREEN,True); text(slide,(859,404,1487,436),"Power Sampling / Self-Consistency",13,MUTED,font_name=FONT_LATIN)
    connector(slide,859,460,1487,460,"C7D9CF",1.0)
    text(slide,(859,496,922,530),"问题",12,GREEN,True); text(slide,(939,496,1487,536),"轨迹质量提高，不等于多轨迹决策更可靠",14,INK_2)
    text(slide,(859,552,922,586),"想法",12,GREEN,True); text(slide,(939,552,1487,592),"Relative-Rank SoftSat 校准锐化强度",14,INK_2)
    text(slide,(859,618,1028,672),"7.42×",29,GREEN,True,font_name=FONT_LATIN); text(slide,(1040,638,1250,668),"加权推理加速",11,MUTED)
    pill(slide,(1323,662,1487,712),"GitHub  ↗",WHITE,GREEN,12,True,"https://github.com/jiaxingsunpku/SoftSat")
    footer(slide,"本文件完整展开项目一，并在项目二正文开始前结束。")

def slide04(slide):
    rect(slide,(0,0,W,H),PLUM,None,False); rect(slide,(0,0,18,H),ROSE,None,False)
    text(slide,(112,195,1496,245),"RESEARCH 01",14,LINE,True,font_name=FONT_LATIN)
    text(slide,(112,288,1496,370),"熵自适应策略蒸馏",42,WHITE,True)
    text(slide,(112,405,1496,455),"为每条 Agent trajectory 找到自己的可蒸馏边界",20,PANEL_2)
    rect(slide,(104,548,1496,706),"654052",PLUM_LIGHT,True,1.0); text(slide,(140,580,300,612),"核心问题",12,LINE)
    text(slide,(140,624,1450,674),"Teacher 的监督并非在整条交互轨迹上都同样可靠。",21,WHITE,True)
    text(slide,(1450,842,1542,872),"04 / 10",10,LINE,False,PP_ALIGN.RIGHT,font_name=FONT_LATIN)

def slide05(slide):
    base(slide,"问题背景：Agentic OPD 会把局部偏差放大成状态漂移",5)
    text(slide,(58,140,1542,180),"在推理题里，错误主要污染后续 token；在交互环境里，动作还会改变下一时刻的世界状态。",15,INK_2)
    xs=[84,382,680,978,1276]; items=[("Student action","早期动作偏差"),("Environment","状态发生改变"),("Observation","后续上下文漂移"),("Teacher","进入陌生支持集"),("Dense loss","监督噪声累积")]
    for i,(x,(a,b)) in enumerate(zip(xs,items),1):
        accent=PLUM if i<4 else ROSE; rect(slide,(x,250,x+224,465),PANEL if i<4 else PANEL_2,LINE,True)
        ellipse(slide,(x+78,277,x+146,345),accent,None); text(slide,(x+78,277,x+146,345),str(i),16,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE,FONT_LATIN)
        text(slide,(x+15,370,x+209,404),a,12,accent,True,PP_ALIGN.CENTER,font_name=FONT_LATIN); text(slide,(x+15,418,x+209,452),b,12,INK_2,False,PP_ALIGN.CENTER)
        if i<5: arrow(slide,x+235,357,x+277,PLUM_LIGHT)
    rect(slide,(58,535,1048,775),PLUM,None,True); text(slide,(92,566,300,602),"为什么重要？",16,WHITE,True)
    bullet(slide,94,620,910,"长程任务训练成本高，后半段低质量 dense supervision 会反向伤害 Student。",14,WHITE,LINE)
    bullet(slide,94,675,910,"统一截断会错杀仍然可教的轨迹，也会放过已经漂移的轨迹。",14,WHITE,LINE)
    bullet(slide,94,730,910,"真正需要识别的是 trajectory-level distillability，而不是固定 turn 数。",14,WHITE,LINE)
    rect(slide,(1090,535,1542,775),GREEN_PANEL,"C7D9CF",True); text(slide,(1124,566,1488,604),"研究目标",15,GREEN,True)
    text(slide,(1124,622,1495,735),"用在线可观测信号，\n为每条轨迹动态决定\n哪些 turn 进入训练 loss。",18,INK_2,True)

def slide06(slide):
    base(slide,"已有方法：TCOD 提升了水平，但仍是全局固定 curriculum",6)
    text(slide,(58,140,1542,180),"为保证可比，下面只报告同一学生、教师、训练步数和 full274 evaluator 下的本地结果。",14,INK_2)
    rect(slide,(58,205,755,700),PANEL,LINE,True); pill(slide,(92,235,275,285),"Vanilla OPD",PLUM,WHITE,12,True)
    text(slide,(92,320,700,360),"全轨迹蒸馏",19,INK,True); text(slide,(92,372,700,404),"所有 turn 都进入 dense loss",13,MUTED)
    for j in range(3): trajectory(slide,92,445+j*66,f"traj {j+1}",6,PLUM_LIGHT)
    text(slide,(92,630,290,680),"79.56%",27,ROSE,True,font_name=FONT_LATIN); text(slide,(300,648,430,675),"218 / 274",11,MUTED,font_name=FONT_LATIN)
    rect(slide,(805,205,1542,700),GREEN_PANEL,"C7D9CF",True); pill(slide,(839,235,1012,285),"TCOD-F2B",GREEN,WHITE,12,True)
    text(slide,(839,320,1480,360),"按训练进度逐步扩展 horizon",19,INK,True); text(slide,(839,372,1480,404),"同一阶段，所有轨迹使用同一个 K",13,MUTED)
    for j in range(3): trajectory(slide,839,445+j*66,f"traj {j+1}",3,GREEN)
    text(slide,(839,630,1037,680),"84.67%",27,GREEN,True,font_name=FONT_LATIN); text(slide,(1047,648,1180,675),"232 / 274",11,MUTED,font_name=FONT_LATIN)
    rect(slide,(58,732,1542,820),PLUM,None,True); text(slide,(92,754,160,798),"局限",14,LINE,True, valign=MSO_ANCHOR.MIDDLE)
    text(slide,(182,748,1490,804),"global schedule 不知道哪条轨迹已经漂移，也不知道哪条轨迹仍值得继续蒸馏。",16,WHITE,True,valign=MSO_ANCHOR.MIDDLE)
    footer(slide,"TCOD 论文：arXiv:2604.24005；本地复现协议见项目 BASELINE_SPEC。")
    # Clickable source label area.
    s=rect(slide,(58,828,650,874),None,None,False); s.click_action.hyperlink.address="https://arxiv.org/abs/2604.24005"

def slide07(slide):
    base(slide,"我们的改进：用教师熵的相对漂移定位 Distillable Frontier",7)
    text(slide,(58,140,1542,180),"不是问“第几轮一定不可靠”，而是问“这条轨迹相对它自己的起点，何时发生持续漂移”。",15,INK_2)
    rect(slide,(58,210,900,690),PANEL,LINE,True); text(slide,(92,238,430,270),"Teacher entropy drift",13,PLUM,True,font_name=FONT_LATIN)
    x0,y0,x1,y1=118,615,846,310; connector(slide,x0,y0,x1,y0,MUTED,1.2); connector(slide,x0,y0,x0,y1,MUTED,1.2)
    vals=[.18,.20,.19,.23,.25,.29,.37,.49,.58,.62,.67]; pts=[]
    for i,v in enumerate(vals): pts.append((x0+30+i*65,y0-28-v*360))
    for a,b in zip(pts,pts[1:]): connector(slide,*a,*b,PLUM,3.0)
    for x,y in pts: ellipse(slide,(x-5,y-5,x+5,y+5),PLUM,None)
    fx=pts[7][0]; connector(slide,fx,y1+15,fx,y0,ROSE,2.2); text(slide,(fx+10,y1+12,fx+160,y1+42),"frontier  fᵢ",12,ROSE,True,font_name=FONT_LATIN)
    rect(slide,(x0,y0+18,fx-4,y0+55),GREEN,None,False); rect(slide,(fx+4,y0+18,x1,y0+55),LINE,None,False)
    text(slide,(x0,y0+18,fx-4,y0+55),"保留 loss",12,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    text(slide,(fx+4,y0+18,x1,y0+55),"屏蔽 suffix",12,MUTED,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    rect(slide,(940,210,1542,690),WHITE,PLUM,True,1.8); text(slide,(978,242,1488,282),"算法定义",17,PLUM,True)
    text(slide,(978,302,1490,334),"1. 完整 rollout，并由 Teacher 评分",14,INK_2)
    text(slide,(978,358,1490,390),"2. 首 3 turn 建立局部熵基线",14,INK_2)
    text(slide,(1000,410,1490,448),"Bᵢ = ⅓ Σₜ₌₀² Hᵢ,ₜ",19,PLUM,True,font_name=FONT_MATH)
    text(slide,(978,476,1490,508),"3. 连续 3 turn 的平均漂移 ≥ τ",14,INK_2)
    text(slide,(1000,528,1490,566),"fᵢ = first sustained crossing",18,ROSE,True,font_name=FONT_MATH)
    text(slide,(978,588,1490,620),"4. 仅保留 t < fᵢ 的 OPD loss",14,INK_2); text(slide,(978,635,1490,664),"无触发则保留整条轨迹",12,MUTED)
    rect(slide,(58,728,1542,816),GREEN_PANEL,"C7D9CF",True); text(slide,(92,749,205,798),"关键控制",14,GREEN,True,valign=MSO_ANCHOR.MIDDLE)
    text(slide,(228,745,1490,802),"环境交互和 Teacher 评分仍跑完整轨迹，只改变 loss selection，因此提升可归因于数据选择。",15,INK_2,True,valign=MSO_ANCHOR.MIDDLE)
    footer(slide,"实现信号：Teacher top-16 partial entropy；τ 为相对熵漂移阈值。")

def slide08(slide):
    base(slide,"实验设计：冻结所有变量，只扫描熵阈值 τ",8)
    info=[("Student","Qwen2.5-3B",PLUM,PANEL),("Teacher","GiGPO 7B",GREEN,GREEN_PANEL),("Training","250 steps · 4 GPU",ROSE,PANEL),("Evaluation","ALFWorld full274",PLUM,PANEL)]
    for i,(k,v,c,f) in enumerate(info):
        x=58+i*371; rect(slide,(x,145,x+334,262),f,LINE,True); text(slide,(x+24,166,x+300,192),k,10,c,True,font_name=FONT_LATIN); text(slide,(x+24,208,x+310,240),v,15,INK,True,font_name=FONT_LATIN)
    rect(slide,(58,305,1050,785),WHITE,LINE,True); text(slide,(92,335,1015,372),"阈值扫描：成功率并非单调，τ = 0.10 形成最佳平衡",16,INK,True)
    cx0,cy0,cx1,cy1=132,700,980,405
    connector(slide,cx0,cy0,cx1,cy0,MUTED,1.1); connector(slide,cx0,cy0,cx0,cy1,MUTED,1.1)
    for v in [78,82,86,90]:
        y=cy0-(v-78)/12*(cy0-cy1); connector(slide,cx0,y,cx1,y,LINE,0.7); text(slide,(74,y-10,cx0-15,y+14),f"{v}%",9,MUTED,False,PP_ALIGN.RIGHT,font_name=FONT_LATIN)
    taus=[.05,.075,.10,.125]; rates=[81.75,80.29,86.86,82.48]; pts=[]
    for i,(t,r) in enumerate(zip(taus,rates)):
        x=cx0+85+i*235; y=cy0-(r-78)/12*(cy0-cy1); pts.append((x,y))
    for a,b in zip(pts,pts[1:]): connector(slide,*a,*b,PLUM_LIGHT,2.6)
    for (x,y),t,r in zip(pts,taus,rates):
        c=ROSE if abs(t-.1)<1e-8 else PLUM; ellipse(slide,(x-8,y-8,x+8,y+8),c,None)
        text(slide,(x-65,y-38,x+65,y-12),f"{r:.2f}%",12,c,True,PP_ALIGN.CENTER,font_name=FONT_LATIN)
        text(slide,(x-65,cy0+20,x+65,cy0+46),f"τ={t:g}",10,MUTED,False,PP_ALIGN.CENTER,font_name=FONT_LATIN)
    rect(slide,(1090,305,1542,785),PLUM,None,True); text(slide,(1124,338,1500,380),"为什么不是越低越好？",16,WHITE,True)
    entries=[("τ 太小","过度截断，丢失仍可学习的后续状态"),("τ 太大","过滤太弱，漂移 suffix 重新进入 loss"),("τ = 0.10","在过滤强度与状态覆盖之间取得平衡")]
    for i,(a,b) in enumerate(entries):
        y=418+i*112; text(slide,(1124,y,1495,y+30),a,14,LINE,True); text(slide,(1124,y+40,1495,y+86),b,13,WHITE)
    footer(slide,"四个阈值均为 fresh 250-step、4-GPU、seed 42 训练；评测协议完全一致。")

def slide09(slide):
    base(slide,"改进效果：τ = 0.10 在最终 full274 取得最高观测成绩",9)
    rect(slide,(58,150,1035,780),PANEL,LINE,True); text(slide,(92,178,500,210),"ALFWorld success rate",13,PLUM,True,font_name=FONT_LATIN)
    x0,y0=142,700; connector(slide,x0,y0,960,y0,MUTED,1.1)
    methods=["Vanilla OPD","TCOD-F2B","Entropy Adaptive"]; rates=[79.56,84.67,86.86]; colors=[PLUM_LIGHT,GREEN,ROSE]
    for i,(m,r,c) in enumerate(zip(methods,rates,colors)):
        x=205+i*255; bh=(r-70)/20*390; rect(slide,(x,y0-bh,x+150,y0),c,None,True)
        text(slide,(x-20,y0-bh-40,x+170,y0-bh-8),f"{r:.2f}%",18,c,True,PP_ALIGN.CENTER,font_name=FONT_LATIN)
        text(slide,(x-35,y0+25,x+185,y0+55),m,10,INK_2,True,PP_ALIGN.CENTER,font_name=FONT_LATIN)
    for v,y in [(70,700),(80,505),(90,310)]:
        connector(slide,142,y,960,y,LINE,0.7); text(slide,(90,y-10,130,y+12),str(v),9,MUTED,False,PP_ALIGN.RIGHT,font_name=FONT_LATIN)
    rect(slide,(1080,150,1542,348),PLUM,None,True); text(slide,(1112,180,1500,210),"对 Vanilla OPD",12,LINE)
    text(slide,(1112,224,1500,272),"+7.30 pp",28,WHITE,True,font_name=FONT_LATIN); text(slide,(1112,292,1500,320),"238 vs 218 / 274",11,PANEL_2,font_name=FONT_LATIN)
    rect(slide,(1080,375,1542,573),GREEN_PANEL,"C7D9CF",True); text(slide,(1112,405,1500,435),"对 TCOD-F2B",12,GREEN)
    text(slide,(1112,449,1500,497),"+2.19 pp",28,GREEN,True,font_name=FONT_LATIN); text(slide,(1112,517,1500,545),"238 vs 232 / 274",11,MUTED,font_name=FONT_LATIN)
    rect(slide,(1080,600,1542,780),PANEL_2,LINE,True); text(slide,(1112,627,1500,657),"泛化拆分",13,PLUM,True)
    text(slide,(1112,678,1215,706),"Seen",11,MUTED,font_name=FONT_LATIN); text(slide,(1242,671,1485,710),"87.86%",18,PLUM,True,font_name=FONT_LATIN)
    text(slide,(1112,731,1215,759),"Unseen",11,MUTED,font_name=FONT_LATIN); text(slide,(1242,724,1485,763),"85.82%",18,PLUM,True,font_name=FONT_LATIN)
    footer(slide,"同协议：140 Seen + 134 Unseen，horizon 30，strict action parser，seed 42。")

def slide10(slide):
    base(slide,"为什么能工作：模型越成熟，算法越少干预",10)
    text(slide,(58,140,1542,180),"τ = 0.10 学到的不是一条死板 schedule：早期积极过滤，后期自然接近 full-loss。",15,INK_2)
    labels=["0–20","20–40","40–60","60–80","80–100"]
    mini_line(slide,(58,215,775,555),[48.02,36.22,21.08,11.35,9.01],ROSE,labels,"Frontier 触发率：48.0% → 9.0%")
    mini_line(slide,(825,215,1542,555),[20.01,22.13,25.66,27.60,28.23],GREEN,labels,"平均可蒸馏 horizon：20.0 → 28.2","")
    rect(slide,(58,592,790,810),PLUM,None,True); text(slide,(92,620,300,654),"我们的结论",14,LINE,True)
    text(slide,(92,670,740,780),"相对熵漂移是一个有效的 trajectory-level\n筛选信号；关键不在“熵越低越好”，而在\n过滤强度与长程状态覆盖之间的平衡。",16,WHITE,True)
    rect(slide,(835,592,1542,810),PANEL_2,LINE,True); text(slide,(869,620,1490,654),"证据边界 / 下一步",14,PLUM,True)
    bullet(slide,871,670,620,"优势最强在 final checkpoint，并非全训练阶段持续领先。",12,INK_2,PLUM_LIGHT)
    bullet(slide,871,717,620,"对 TCOD 的 +2.19 pp 在单种子 McNemar 检验下尚不显著（p=.362）。",12,INK_2,PLUM_LIGHT)
    bullet(slide,871,764,620,"当前仍跑完整 rollout；下一步验证多种子稳定性，并探索连续加权与算力分配。",12,INK_2,PLUM_LIGHT)
    footer(slide,"项目一小结 · 下一页将进入 SoftSat，但本次交付按要求在此结束。")

def main():
    OUT_DIR.mkdir(parents=True,exist_ok=True)
    prs=Presentation(str(REFERENCE)); remove_all_slides(prs); prs.slide_width=Inches(13.333333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[0]
    for fn in [slide01,slide02,slide03,slide04,slide05,slide06,slide07,slide08,slide09,slide10]:
        slide=prs.slides.add_slide(blank); fn(slide)
    prs.core_properties.title="孙家兴｜面试自我介绍｜熵自适应策略蒸馏"
    prs.core_properties.subject="Editable interview presentation"
    prs.core_properties.author="孙家兴"
    out=OUT_DIR/"孙家兴_面试自我介绍_OPD项目_可编辑.pptx"; prs.save(out); print(out)
if __name__=='__main__': main()
