#!/usr/bin/env python3
"""Generate a visually frozen 16:9 interview deck and matching speaker notes.

The reference deck is used as the PowerPoint package/theme carrier.  Each slide
is rendered to a high-resolution PNG first, then embedded as a full-slide image.
This keeps the layout stable across machines without requiring the original
authoring fonts.  URL hit-boxes are added on top of the relevant labels.
"""

from __future__ import annotations

import math
import os
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "deliverables" / "interview_self_intro_opd"
PREVIEW_DIR = OUT_DIR / "slides"
REFERENCE = Path(
    "/vepfs-mlp2/mlp-public/252302025/sjx/.codex-sjx/attachments/"
    "3147fb86-839f-4e03-b4c3-0730d2b8a9cd/"
    "on-policy-distillation-agentic-opd-revised.pptx"
)

W, H = 1600, 900
FONT_CJK = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
FONT_LATIN_REG = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_LATIN_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

# Reference-template palette.
INK = "#211F20"
INK_2 = "#343033"
MUTED = "#716A6E"
PLUM = "#70485B"
PLUM_LIGHT = "#9B7184"
ROSE = "#8C3E49"
GREEN = "#315C4A"
PAPER = "#FFFFFF"
PANEL = "#F4F0F2"
PANEL_2 = "#EEE3E8"
LINE = "#D8D0D4"
GREEN_PANEL = "#E5EEE9"
WHITE = "#FFFFFF"


def font(size: int, bold: bool = False, latin: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_CJK
    if latin:
        path = FONT_LATIN_BOLD if bold else FONT_LATIN_REG
    return ImageFont.truetype(path, size=size)


def is_cjk(ch: str) -> bool:
    code = ord(ch)
    return (
        0x3400 <= code <= 0x4DBF
        or 0x4E00 <= code <= 0x9FFF
        or 0xF900 <= code <= 0xFAFF
        or 0x3000 <= code <= 0x303F
        or 0xFF00 <= code <= 0xFFEF
    )


def glyph_font(ch: str, size: int, bold: bool, force_latin: bool = False) -> ImageFont.FreeTypeFont:
    # DroidSansFallback is a CJK-only fallback on this host. Mixing it with
    # DejaVu per character prevents Latin words, digits, Greek symbols, and
    # mathematical punctuation from turning into tofu boxes.
    return font(size, bold=bold, latin=force_latin or not is_cjk(ch))


def mixed_line_width(text: str, size: int, bold: bool = False, force_latin: bool = False) -> int:
    width = 0.0
    for ch in text:
        f = glyph_font(ch, size, bold, force_latin)
        width += f.getlength(ch)
    return int(math.ceil(width))


def draw_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    size: int,
    color: str = INK,
    bold: bool = False,
    anchor: str | None = None,
    latin: bool = False,
    spacing: int = 8,
) -> None:
    lines = text.split("\n")
    ascent, descent = font(size, bold=bold, latin=True).getmetrics()
    line_height = ascent + descent + spacing
    block_height = (len(lines) - 1) * line_height + ascent + descent
    h_align = "l"
    v_align = "t"
    if anchor:
        h_align = anchor[0]
        v_align = anchor[1]

    x_anchor, y_anchor = xy
    if v_align == "m":
        y_start = y_anchor - block_height / 2
    elif v_align in ("s", "b", "d"):
        y_start = y_anchor - ascent
    else:
        y_start = y_anchor

    for line_idx, line in enumerate(lines):
        line_width = mixed_line_width(line, size, bold, latin)
        if h_align == "m":
            x = x_anchor - line_width / 2
        elif h_align == "r":
            x = x_anchor - line_width
        else:
            x = x_anchor
        y = y_start + line_idx * line_height
        for ch in line:
            f = glyph_font(ch, size, bold, latin)
            kwargs = {"fill": color, "font": f, "anchor": "la"}
            if bold and is_cjk(ch) and not latin:
                kwargs["stroke_width"] = 1
                kwargs["stroke_fill"] = color
            draw.text((int(x), int(y)), ch, **kwargs)
            x += f.getlength(ch)


def wrap(draw: ImageDraw.ImageDraw, text: str, max_width: int, size: int, bold: bool = False) -> str:
    lines: list[str] = []
    for raw in text.split("\n"):
        cur = ""
        for ch in raw:
            trial = cur + ch
            if mixed_line_width(trial, size, bold) <= max_width:
                cur = trial
            else:
                if cur:
                    lines.append(cur)
                cur = ch
        lines.append(cur)
    return "\n".join(lines)


def rounded(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill: str, radius: int = 24,
            outline: str | None = None, width: int = 2) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def base_slide(title: str, number: int, total: int = 10, kicker: str | None = None) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    im = Image.new("RGB", (W, H), PAPER)
    d = ImageDraw.Draw(im)
    draw_text(d, (58, 32), title, 39, PLUM, bold=True)
    if kicker:
        draw_text(d, (58, 90), kicker, 18, MUTED)
    d.line((58, 114, 1542, 114), fill=LINE, width=2)
    draw_text(d, (1542, 850), f"{number:02d} / {total:02d}", 15, MUTED, anchor="ra", latin=True)
    return im, d


def footer(d: ImageDraw.ImageDraw, text: str) -> None:
    draw_text(d, (58, 850), text, 16, MUTED)


def pill(d: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str, fill: str, color: str = INK,
         size: int = 18, bold: bool = False) -> None:
    rounded(d, box, fill, radius=18)
    x1, y1, x2, y2 = box
    draw_text(d, ((x1 + x2) // 2, (y1 + y2) // 2), text, size, color, bold=bold, anchor="mm")


def card(d: ImageDraw.ImageDraw, box: tuple[int, int, int, int], index: str, title: str, body: str,
         accent: str = PLUM, fill: str = PANEL) -> None:
    rounded(d, box, fill, radius=26, outline=LINE)
    x1, y1, x2, y2 = box
    draw_text(d, (x1 + 28, y1 + 24), index, 19, accent, bold=True, latin=True)
    draw_text(d, (x1 + 28, y1 + 65), title, 26, INK, bold=True)
    wrapped = wrap(d, body, x2 - x1 - 56, 19)
    draw_text(d, (x1 + 28, y1 + 118), wrapped, 19, MUTED, spacing=10)


def slide_01() -> Image.Image:
    im = Image.new("RGB", (W, H), PAPER)
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 38, H), fill=PLUM)
    d.rectangle((1220, 0, W, H), fill=PANEL)
    d.ellipse((1300, 95, 1515, 310), fill=PLUM)
    draw_text(d, (1408, 203), "SJX", 54, WHITE, bold=True, anchor="mm", latin=True)
    draw_text(d, (98, 178), "孙家兴", 70, PLUM, bold=True)
    draw_text(d, (98, 286), "面试自我介绍", 43, INK, bold=True)
    draw_text(d, (98, 380), "大模型后训练  ·  Agentic Learning  ·  推理增强", 25, MUTED)
    d.line((98, 478, 1050, 478), fill=LINE, width=3)
    draw_text(d, (98, 525), "北京大学 · 智能学院", 25, INK_2, bold=True)
    draw_text(d, (98, 585), "M.S. Candidate  |  August 2026", 20, MUTED, latin=True)
    pill(d, (1270, 400, 1550, 464), "Research × Engineering", WHITE, PLUM, 18, True)
    draw_text(d, (1270, 520), "从问题定义、可复现基线，\n到机制诊断与完整评测", 22, INK_2, spacing=13)
    draw_text(d, (98, 820), "2501213407@stu.pku.edu.cn  ·  +86 188-1171-7335", 17, MUTED)
    draw_text(d, (1542, 850), "01 / 10", 15, MUTED, anchor="ra", latin=True)
    return im


def slide_02() -> Image.Image:
    im, d = base_slide("基本信息：研究型算法工程与可验证实验", 2)
    rounded(d, (58, 155, 505, 800), PLUM, radius=28)
    draw_text(d, (94, 195), "孙家兴", 48, WHITE, bold=True)
    draw_text(d, (94, 270), "北京大学 · 智能学院", 25, WHITE, bold=True)
    draw_text(d, (94, 315), "硕士研究生（2025.09—至今）", 19, "#EEE3E8")
    d.line((94, 372, 465, 372), fill=PLUM_LIGHT, width=2)
    draw_text(d, (94, 410), "教育背景", 18, "#D8D0D4")
    draw_text(d, (94, 452), "北京大学 · 信息科学技术学院\n本科（2021.09—2025.06）", 21, WHITE, spacing=12)
    draw_text(d, (94, 570), "联系方式", 18, "#D8D0D4")
    draw_text(d, (94, 612), "188-1171-7335\n2501213407@stu.pku.edu.cn", 19, WHITE, spacing=12)
    pill(d, (94, 720, 318, 770), "GitHub  ↗", WHITE, PLUM, 18, True)

    card(d, (550, 155, 1015, 440), "01", "研究方向",
         "在线策略蒸馏 · 长程 Agent\n采样与多轨迹决策 · 自动评测", PLUM, PANEL)
    card(d, (1050, 155, 1542, 440), "02", "训练与系统",
         "PyTorch · veRL · vLLM · FSDP\nRay · Tool Calling · JSONL artifacts", GREEN, GREEN_PANEL)
    card(d, (550, 475, 1015, 800), "03", "我的工作方式",
         "先冻结可比较的基线，再做诊断；\n把每个提升拆成可证伪的假设。", ROSE, PANEL_2)
    card(d, (1050, 475, 1542, 800), "04", "相关经历",
         "大模型蒸馏、代码 Agent、\nPPO / MAPPO 多智能体控制。", PLUM, PANEL)
    return im


def slide_03() -> Image.Image:
    im, d = base_slide("两项核心研究：从数据机制到决策机制", 3)
    draw_text(d, (58, 145), "共同主线：不是只追求一个更高分，而是解释模型为什么失效、什么信号能修复。", 23, INK_2)

    rounded(d, (58, 220, 775, 742), PANEL, radius=28, outline=PLUM, width=3)
    pill(d, (92, 252, 255, 302), "PROJECT 01", PLUM, WHITE, 16, True)
    draw_text(d, (92, 345), "熵自适应策略蒸馏", 36, PLUM, bold=True)
    draw_text(d, (92, 407), "Agentic OPD / ALFWorld", 20, MUTED, latin=True)
    d.line((92, 460, 720, 460), fill=LINE, width=2)
    draw_text(d, (92, 500), "问题", 18, PLUM, bold=True)
    draw_text(d, (172, 500), "固定 horizon 无法判断每条轨迹何时开始漂移", 20, INK_2)
    draw_text(d, (92, 555), "想法", 18, PLUM, bold=True)
    draw_text(d, (172, 555), "以教师熵的相对变化定位 distillable frontier", 20, INK_2)
    draw_text(d, (92, 625), "86.86%", 43, ROSE, bold=True, latin=True)
    draw_text(d, (300, 640), "ALFWorld full274", 18, MUTED, latin=True)
    pill(d, (556, 662, 720, 712), "GitHub  ↗", WHITE, PLUM, 17, True)

    rounded(d, (825, 220, 1542, 742), GREEN_PANEL, radius=28, outline=GREEN, width=3)
    pill(d, (859, 252, 1022, 302), "PROJECT 02", GREEN, WHITE, 16, True)
    draw_text(d, (859, 345), "SoftSat：采样失效修复", 35, GREEN, bold=True)
    draw_text(d, (859, 407), "Power Sampling / Self-Consistency", 20, MUTED, latin=True)
    d.line((859, 460, 1487, 460), fill="#C7D9CF", width=2)
    draw_text(d, (859, 500), "问题", 18, GREEN, bold=True)
    draw_text(d, (939, 500), "轨迹质量提高，不等于多轨迹决策更可靠", 20, INK_2)
    draw_text(d, (859, 555), "想法", 18, GREEN, bold=True)
    draw_text(d, (939, 555), "Relative-Rank SoftSat 校准锐化强度", 20, INK_2)
    draw_text(d, (859, 625), "7.42×", 43, GREEN, bold=True, latin=True)
    draw_text(d, (1040, 640), "加权推理加速", 18, MUTED)
    pill(d, (1323, 662, 1487, 712), "GitHub  ↗", WHITE, GREEN, 17, True)
    footer(d, "本文件完整展开项目一，并在项目二正文开始前结束。")
    return im


def slide_04() -> Image.Image:
    im = Image.new("RGB", (W, H), PLUM)
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 18, H), fill=ROSE)
    draw_text(d, (112, 205), "RESEARCH 01", 22, "#D8D0D4", bold=True, latin=True)
    draw_text(d, (112, 295), "熵自适应策略蒸馏", 62, WHITE, bold=True)
    draw_text(d, (112, 405), "为每条 Agent trajectory 找到自己的可蒸馏边界", 29, "#EEE3E8")
    rounded(d, (104, 548, 1496, 706), "#654052", radius=24, outline=PLUM_LIGHT)
    draw_text(d, (140, 580), "核心问题", 18, "#D8D0D4")
    draw_text(d, (140, 625), "Teacher 的监督并非在整条交互轨迹上都同样可靠。", 30, WHITE, bold=True)
    draw_text(d, (1542, 850), "04 / 10", 15, "#D8D0D4", anchor="ra", latin=True)
    return im


def slide_05() -> Image.Image:
    im, d = base_slide("问题背景：Agentic OPD 会把局部偏差放大成状态漂移", 5)
    draw_text(d, (58, 145), "在推理题里，错误主要污染后续 token；在交互环境里，动作还会改变下一时刻的世界状态。", 22, INK_2)

    xs = [84, 382, 680, 978, 1276]
    labels = [
        ("Student action", "早期动作偏差"),
        ("Environment", "状态发生改变"),
        ("Observation", "后续上下文漂移"),
        ("Teacher", "进入陌生支持集"),
        ("Dense loss", "监督噪声累积"),
    ]
    for i, (x, (a, b)) in enumerate(zip(xs, labels), start=1):
        fill = PANEL if i < 4 else PANEL_2
        accent = PLUM if i < 4 else ROSE
        rounded(d, (x, 250, x + 224, 465), fill, radius=24, outline=LINE)
        d.ellipse((x + 78, 277, x + 146, 345), fill=accent)
        draw_text(d, (x + 112, 311), str(i), 24, WHITE, bold=True, anchor="mm", latin=True)
        draw_text(d, (x + 112, 374), a, 18, accent, bold=True, anchor="mm", latin=True)
        draw_text(d, (x + 112, 423), b, 18, INK_2, anchor="mm")
        if i < 5:
            d.line((x + 235, 357, x + 277, 357), fill=PLUM_LIGHT, width=5)
            d.polygon([(x + 277, 357), (x + 260, 347), (x + 260, 367)], fill=PLUM_LIGHT)

    rounded(d, (58, 535, 1048, 775), PLUM, radius=26)
    draw_text(d, (92, 570), "为什么重要？", 24, WHITE, bold=True)
    bullets = [
        "长程任务的训练成本高，后半段的低质量 dense supervision 会反向伤害 Student。",
        "统一截断会错杀仍然可教的轨迹，也会放过已经漂移的轨迹。",
        "真正需要识别的是 trajectory-level distillability，而不是固定 turn 数。",
    ]
    for i, b in enumerate(bullets):
        y = 625 + i * 52
        d.ellipse((94, y + 7, 106, y + 19), fill="#D8D0D4")
        draw_text(d, (124, y), b, 20, WHITE)

    rounded(d, (1090, 535, 1542, 775), GREEN_PANEL, radius=26, outline="#C7D9CF")
    draw_text(d, (1124, 570), "研究目标", 22, GREEN, bold=True)
    draw_text(d, (1124, 625), "用在线可观测信号，\n为每条轨迹动态决定\n哪些 turn 进入训练 loss。", 26, INK_2, bold=True, spacing=14)
    return im


def trajectory_row(d: ImageDraw.ImageDraw, x: int, y: int, widths: list[int], colors: list[str], label: str) -> None:
    draw_text(d, (x, y - 8), label, 18, MUTED)
    px = x + 130
    for w, c in zip(widths, colors):
        rounded(d, (px, y, px + w, y + 38), c, radius=8)
        px += w + 8


def slide_06() -> Image.Image:
    im, d = base_slide("已有方法：TCOD 提升了水平，但仍是全局固定 curriculum", 6)
    draw_text(d, (58, 145), "为保证可比，下面只报告同一学生、教师、训练步数和 full274 evaluator 下的本地结果。", 21, INK_2)

    rounded(d, (58, 205, 755, 700), PANEL, radius=26, outline=LINE)
    pill(d, (92, 235, 275, 285), "Vanilla OPD", PLUM, WHITE, 17, True)
    draw_text(d, (92, 325), "全轨迹蒸馏", 28, INK, bold=True)
    draw_text(d, (92, 375), "所有 turn 都进入 dense loss", 19, MUTED)
    for j in range(3):
        trajectory_row(d, 92, 445 + j * 66, [72] * 6, [PLUM_LIGHT] * 6, f"traj {j+1}")
    draw_text(d, (92, 640), "79.56%", 38, ROSE, bold=True, latin=True)
    draw_text(d, (300, 650), "218 / 274", 18, MUTED, latin=True)

    rounded(d, (805, 205, 1542, 700), GREEN_PANEL, radius=26, outline="#C7D9CF")
    pill(d, (839, 235, 1012, 285), "TCOD-F2B", GREEN, WHITE, 17, True)
    draw_text(d, (839, 325), "按训练进度逐步扩展 horizon", 28, INK, bold=True)
    draw_text(d, (839, 375), "同一阶段，所有轨迹使用同一个 K", 19, MUTED)
    colors = [GREEN, GREEN, GREEN, LINE, LINE, LINE]
    for j in range(3):
        trajectory_row(d, 839, 445 + j * 66, [72] * 6, colors, f"traj {j+1}")
    draw_text(d, (839, 640), "84.67%", 38, GREEN, bold=True, latin=True)
    draw_text(d, (1047, 650), "232 / 274", 18, MUTED, latin=True)

    rounded(d, (58, 732, 1542, 820), PLUM, radius=22)
    draw_text(d, (92, 776), "局限", 20, "#D8D0D4", bold=True, anchor="lm")
    draw_text(d, (182, 776), "global schedule 不知道哪条轨迹已经漂移，也不知道哪条轨迹仍值得继续蒸馏。", 24, WHITE, bold=True, anchor="lm")
    footer(d, "TCOD 论文：arXiv:2604.24005；本地复现协议见项目 BASELINE_SPEC。")
    return im


def slide_07() -> Image.Image:
    im, d = base_slide("我们的改进：用教师熵的相对漂移定位 Distillable Frontier", 7)
    draw_text(d, (58, 145), "不是问“第几轮一定不可靠”，而是问“这条轨迹相对它自己的起点，何时发生持续漂移”。", 22, INK_2)

    # Entropy schematic.
    rounded(d, (58, 210, 900, 690), PANEL, radius=26, outline=LINE)
    draw_text(d, (92, 240), "Teacher entropy drift", 19, PLUM, bold=True, latin=True)
    x0, y0, x1, y1 = 118, 615, 846, 310
    d.line((x0, y0, x1, y0), fill=MUTED, width=2)
    d.line((x0, y0, x0, y1), fill=MUTED, width=2)
    pts = []
    vals = [0.18, 0.20, 0.19, 0.23, 0.25, 0.29, 0.37, 0.49, 0.58, 0.62, 0.67]
    for i, v in enumerate(vals):
        x = x0 + 30 + i * 65
        y = y0 - 28 - v * 360
        pts.append((x, y))
    d.line(pts, fill=PLUM, width=6, joint="curve")
    for x, y in pts:
        d.ellipse((x - 6, y - 6, x + 6, y + 6), fill=PLUM)
    frontier_x = pts[7][0]
    d.line((frontier_x, y1 + 15, frontier_x, y0), fill=ROSE, width=4)
    draw_text(d, (frontier_x + 10, y1 + 12), "frontier  fᵢ", 18, ROSE, bold=True, latin=True)
    d.rectangle((x0 + 1, y0 + 18, frontier_x - 4, y0 + 55), fill=GREEN)
    d.rectangle((frontier_x + 4, y0 + 18, x1, y0 + 55), fill=LINE)
    draw_text(d, ((x0 + frontier_x) // 2, y0 + 36), "保留 loss", 18, WHITE, bold=True, anchor="mm")
    draw_text(d, ((frontier_x + x1) // 2, y0 + 36), "屏蔽 suffix", 18, MUTED, bold=True, anchor="mm")

    rounded(d, (940, 210, 1542, 690), WHITE, radius=26, outline=PLUM, width=3)
    draw_text(d, (978, 245), "算法定义", 24, PLUM, bold=True)
    draw_text(d, (978, 305), "1. 完整 rollout，并由 Teacher 评分", 21, INK_2)
    draw_text(d, (978, 360), "2. 首 3 turn 建立局部熵基线", 21, INK_2)
    draw_text(d, (1000, 414), "Bᵢ = ⅓ Σₜ₌₀² Hᵢ,ₜ", 27, PLUM, bold=True, latin=True)
    draw_text(d, (978, 478), "3. 连续 3 turn 的平均漂移 ≥ τ", 21, INK_2)
    draw_text(d, (1000, 532), "fᵢ = first sustained crossing", 25, ROSE, bold=True, latin=True)
    draw_text(d, (978, 592), "4. 仅保留 t < fᵢ 的 OPD loss", 21, INK_2)
    draw_text(d, (978, 637), "无触发则保留整条轨迹", 18, MUTED)

    rounded(d, (58, 728, 1542, 816), GREEN_PANEL, radius=22, outline="#C7D9CF")
    draw_text(d, (92, 772), "关键控制", 20, GREEN, bold=True, anchor="lm")
    draw_text(d, (228, 772), "环境交互和 Teacher 评分仍跑完整轨迹，只改变 loss selection，因此提升可归因于数据选择。", 22, INK_2, bold=True, anchor="lm")
    footer(d, "实现信号：Teacher top-16 partial entropy；τ 为相对熵漂移阈值。")
    return im


def slide_08() -> Image.Image:
    im, d = base_slide("实验设计：冻结所有变量，只扫描熵阈值 τ", 8)
    cards = [
        ("Student", "Qwen2.5-3B", PLUM),
        ("Teacher", "GiGPO 7B", GREEN),
        ("Training", "250 steps · 4 GPU", ROSE),
        ("Evaluation", "ALFWorld full274", PLUM),
    ]
    for i, (k, v, c) in enumerate(cards):
        x = 58 + i * 371
        rounded(d, (x, 145, x + 334, 262), PANEL if i != 1 else GREEN_PANEL, radius=22, outline=LINE)
        draw_text(d, (x + 24, 170), k, 15, c, bold=True, latin=True)
        draw_text(d, (x + 24, 210), v, 22, INK, bold=True, latin=True)

    rounded(d, (58, 305, 1050, 785), WHITE, radius=26, outline=LINE)
    draw_text(d, (92, 338), "阈值扫描：成功率并非单调，τ = 0.10 形成最佳平衡", 23, INK, bold=True)
    chart_x0, chart_y0, chart_x1, chart_y1 = 132, 700, 980, 405
    d.line((chart_x0, chart_y0, chart_x1, chart_y0), fill=MUTED, width=2)
    d.line((chart_x0, chart_y0, chart_x0, chart_y1), fill=MUTED, width=2)
    for val in [78, 82, 86, 90]:
        y = chart_y0 - (val - 78) / 12 * (chart_y0 - chart_y1)
        d.line((chart_x0, y, chart_x1, y), fill=LINE, width=1)
        draw_text(d, (chart_x0 - 18, int(y)), f"{val}%", 15, MUTED, anchor="rm", latin=True)
    taus = [0.05, 0.075, 0.10, 0.125]
    rates = [81.75, 80.29, 86.86, 82.48]
    pts = []
    for i, (t, r) in enumerate(zip(taus, rates)):
        x = chart_x0 + 85 + i * 235
        y = chart_y0 - (r - 78) / 12 * (chart_y0 - chart_y1)
        pts.append((x, y))
    d.line(pts, fill=PLUM_LIGHT, width=5)
    for i, ((x, y), t, r) in enumerate(zip(pts, taus, rates)):
        c = ROSE if abs(t - 0.1) < 1e-9 else PLUM
        rad = 11 if c == ROSE else 8
        d.ellipse((x - rad, y - rad, x + rad, y + rad), fill=c)
        draw_text(d, (x, y - 30), f"{r:.2f}%", 18, c, bold=True, anchor="ms", latin=True)
        draw_text(d, (x, chart_y0 + 30), f"τ={t:g}", 16, MUTED, anchor="ms", latin=True)

    rounded(d, (1090, 305, 1542, 785), PLUM, radius=26)
    draw_text(d, (1124, 342), "为什么不是越低越好？", 23, WHITE, bold=True)
    bullets = [
        ("τ 太小", "过度截断，丢失仍可学习的后续状态"),
        ("τ 太大", "过滤太弱，漂移 suffix 重新进入 loss"),
        ("τ = 0.10", "在过滤强度与状态覆盖之间取得平衡"),
    ]
    for i, (a, b) in enumerate(bullets):
        y = 420 + i * 112
        draw_text(d, (1124, y), a, 20, "#D8D0D4", bold=True)
        draw_text(d, (1124, y + 42), wrap(d, b, 360, 19), 19, WHITE, spacing=9)
    footer(d, "四个阈值均为 fresh 250-step、4-GPU、seed 42 训练；评测协议完全一致。")
    return im


def slide_09() -> Image.Image:
    im, d = base_slide("改进效果：τ = 0.10 在最终 full274 取得最高观测成绩", 9)
    rounded(d, (58, 150, 1035, 780), PANEL, radius=26, outline=LINE)
    draw_text(d, (92, 182), "ALFWorld success rate", 20, PLUM, bold=True, latin=True)
    x0, y0 = 142, 700
    d.line((x0, y0, 960, y0), fill=MUTED, width=2)
    methods = ["Vanilla OPD", "TCOD-F2B", "Entropy Adaptive"]
    rates = [79.56, 84.67, 86.86]
    colors = [PLUM_LIGHT, GREEN, ROSE]
    for i, (m, r, c) in enumerate(zip(methods, rates, colors)):
        x = 205 + i * 255
        bh = (r - 70) / 20 * 390
        d.rounded_rectangle((x, y0 - bh, x + 150, y0), radius=15, fill=c)
        draw_text(d, (x + 75, y0 - bh - 26), f"{r:.2f}%", 28, c, bold=True, anchor="ms", latin=True)
        draw_text(d, (x + 75, y0 + 32), m, 16, INK_2, bold=True, anchor="ms", latin=True)
    d.line((142, 505, 960, 505), fill=LINE, width=1)
    d.line((142, 310, 960, 310), fill=LINE, width=1)
    draw_text(d, (128, 700), "70", 14, MUTED, anchor="rm", latin=True)
    draw_text(d, (128, 505), "80", 14, MUTED, anchor="rm", latin=True)
    draw_text(d, (128, 310), "90", 14, MUTED, anchor="rm", latin=True)

    rounded(d, (1080, 150, 1542, 348), PLUM, radius=24)
    draw_text(d, (1112, 182), "对 Vanilla OPD", 18, "#D8D0D4")
    draw_text(d, (1112, 229), "+7.30 pp", 40, WHITE, bold=True, latin=True)
    draw_text(d, (1112, 292), "238 vs 218 / 274", 18, "#EEE3E8", latin=True)

    rounded(d, (1080, 375, 1542, 573), GREEN_PANEL, radius=24, outline="#C7D9CF")
    draw_text(d, (1112, 407), "对 TCOD-F2B", 18, GREEN)
    draw_text(d, (1112, 454), "+2.19 pp", 40, GREEN, bold=True, latin=True)
    draw_text(d, (1112, 517), "238 vs 232 / 274", 18, MUTED, latin=True)

    rounded(d, (1080, 600, 1542, 780), PANEL_2, radius=24, outline=LINE)
    draw_text(d, (1112, 630), "泛化拆分", 18, PLUM, bold=True)
    draw_text(d, (1112, 681), "Seen", 17, MUTED, latin=True)
    draw_text(d, (1242, 681), "87.86%", 28, PLUM, bold=True, latin=True)
    draw_text(d, (1112, 734), "Unseen", 17, MUTED, latin=True)
    draw_text(d, (1242, 734), "85.82%", 28, PLUM, bold=True, latin=True)
    footer(d, "同协议：140 Seen + 134 Unseen，horizon 30，strict action parser，seed 42。")
    return im


def mini_line(d: ImageDraw.ImageDraw, box: tuple[int, int, int, int], values: list[float], color: str,
              labels: list[str], title: str, suffix: str = "%") -> None:
    x1, y1, x2, y2 = box
    rounded(d, box, WHITE, radius=22, outline=LINE)
    draw_text(d, (x1 + 28, y1 + 26), title, 20, INK, bold=True)
    px1, py1, px2, py2 = x1 + 54, y2 - 70, x2 - 34, y1 + 94
    vmin, vmax = min(values), max(values)
    span = max(vmax - vmin, 1e-6)
    pts = []
    for i, v in enumerate(values):
        x = px1 + i * (px2 - px1) / (len(values) - 1)
        y = py1 - (v - vmin) / span * (py1 - py2)
        pts.append((x, y))
    d.line(pts, fill=color, width=5)
    for i, ((x, y), v, lab) in enumerate(zip(pts, values, labels)):
        d.ellipse((x - 7, y - 7, x + 7, y + 7), fill=color)
        if i in (0, len(values) - 1):
            draw_text(d, (int(x), int(y - 18)), f"{v:.1f}{suffix}", 16, color, bold=True, anchor="ms", latin=True)
        draw_text(d, (int(x), py1 + 24), lab, 14, MUTED, anchor="ms", latin=True)


def slide_10() -> Image.Image:
    im, d = base_slide("为什么能工作：模型越成熟，算法越少干预", 10)
    draw_text(d, (58, 145), "τ = 0.10 学到的不是一条死板 schedule：早期积极过滤，后期自然接近 full-loss。", 22, INK_2)
    labels = ["0–20", "20–40", "40–60", "60–80", "80–100"]
    mini_line(d, (58, 215, 775, 555), [48.02, 36.22, 21.08, 11.35, 9.01], ROSE, labels,
              "Frontier 触发率：48.0% → 9.0%")
    mini_line(d, (825, 215, 1542, 555), [20.01, 22.13, 25.66, 27.60, 28.23], GREEN, labels,
              "平均可蒸馏 horizon：20.0 → 28.2", "")

    rounded(d, (58, 592, 790, 810), PLUM, radius=26)
    draw_text(d, (92, 622), "我们的结论", 21, "#D8D0D4", bold=True)
    draw_text(d, (92, 670), "相对熵漂移是一个有效的 trajectory-level\n筛选信号；关键不在“熵越低越好”，而在\n过滤强度与长程状态覆盖之间的平衡。", 23, WHITE, bold=True, spacing=13)

    rounded(d, (835, 592, 1542, 810), PANEL_2, radius=26, outline=LINE)
    draw_text(d, (869, 622), "证据边界 / 下一步", 21, PLUM, bold=True)
    notes = [
        "优势最强在 final checkpoint，并非全训练阶段持续领先。",
        "对 TCOD 的 +2.19 pp 在单种子 McNemar 检验下尚不显著（p=.362）。",
        "当前仍跑完整 rollout；下一步验证多种子稳定性，并探索连续加权与算力分配。",
    ]
    for i, b in enumerate(notes):
        y = 670 + i * 47
        d.ellipse((871, y + 7, 883, y + 19), fill=PLUM_LIGHT)
        draw_text(d, (900, y), wrap(d, b, 605, 17), 17, INK_2, spacing=8)
    footer(d, "项目一小结 · 下一页将进入 SoftSat，但本次交付按要求在此结束。")
    return im


SLIDE_BUILDERS = [
    slide_01,
    slide_02,
    slide_03,
    slide_04,
    slide_05,
    slide_06,
    slide_07,
    slide_08,
    slide_09,
    slide_10,
]


HYPERLINKS = {
    1: [
        ((94, 720, 318, 770), "https://github.com/Sisyphe-lee"),
    ],
    2: [
        ((556, 662, 720, 712), "https://github.com/Sisyphe-lee/opd-baseline-repro"),
        ((1323, 662, 1487, 712), "https://github.com/jiaxingsunpku/SoftSat"),
    ],
    5: [
        ((58, 822, 650, 872), "https://arxiv.org/abs/2604.24005"),
    ],
}


SPEAKER_NOTES = """# 孙家兴｜面试自我介绍 PPT 讲稿

建议总时长：7–9 分钟。以下讲稿不必逐字念，重点是每页只讲一个结论。

## 01 封面（15–20 秒）

大家好，我叫孙家兴，目前是北京大学智能学院硕士研究生。我的研究兴趣集中在大模型后训练、长程智能体学习，以及采样与决策机制。我比较关注的不只是把指标做高，也会把基线、评测协议和机制诊断一起做完整。

## 02 基本信息（35–45 秒）

我本科和硕士都在北京大学。本科阶段做过多智能体强化学习，研究生阶段重点转向大模型蒸馏和 Agent。工程上主要使用 PyTorch、veRL、vLLM、FSDP 和 Ray，也做过代码 Agent 与自动评测。我习惯先冻结可比基线，再把改进拆成可验证的假设。

## 03 两项核心研究（30–40 秒）

目前我最核心的两项研究都围绕“为什么已有方法会失效”。第一项研究长程 Agent 的在线策略蒸馏，用教师熵的动态变化判断每条轨迹何时不再适合蒸馏；第二项研究多轨迹推理里的 Power Sampling，解决单条轨迹质量提高但最终投票变差的问题。接下来先展开第一项。

## 04 项目一过渡（10 秒）

这个项目的核心问题是：Teacher 的监督并不是在整条交互轨迹上都同样可靠，我们能否为每条轨迹找到自己的可蒸馏边界？

## 05 背景与重要性（55–70 秒）

在普通推理任务中，早期错误主要影响后续 token；但在 Agent 环境里，Student 的动作会改变环境状态，新的 observation 又会进入下一轮上下文。因此一个早期偏差会通过 action、transition、observation 被持续放大，最终把 Teacher 带进训练分布之外的状态。此时继续做 dense distillation，后半段 supervision 可能从帮助变成噪声。固定截断也不理想，因为不同轨迹的漂移时刻并不相同。

## 06 已有方法水平与局限（60–75 秒）

Vanilla OPD 对整条轨迹计算蒸馏损失，在我们的同协议 full274 评测中是 79.56%。TCOD 提出按照训练进度逐步扩展 horizon，把结果提升到 84.67%，证明 temporal curriculum 是有效的。但它对同一训练阶段的所有轨迹使用相同的 K，因此无法识别哪条轨迹已经漂移，也无法保留另一条仍然可靠的更长前缀。我们的切入点就是把全局 schedule 改成 trajectory-specific selection。

## 07 我们的改进（80–100 秒）

我们用 Teacher 在每个 turn 的 top-16 partial entropy 作为 probe。先用前三轮建立这条轨迹自己的局部基线，然后检查连续三轮的平均熵漂移是否超过阈值。如果首次持续越过阈值，就把该位置定义为 distillable frontier，只让 frontier 之前的 turn 进入 OPD loss；如果没有触发，就保留完整轨迹。这里最重要的实验控制是：我们仍然跑完整 rollout，也仍然让 Teacher 评分完整轨迹，只改变 loss selection。因此性能变化可以更干净地归因到数据选择，而不是交互预算变化。

## 08 实验设计与阈值（50–65 秒）

学生、教师、训练步数、硬件布局和 evaluator 全部冻结，只扫描阈值。结果不是单调的：阈值太小会过度截断，损失长程状态覆盖；阈值太大则过滤太弱。τ=0.10 在两者之间取得最佳平衡，full274 成功率达到 86.86%。

## 09 改进效果（45–60 秒）

在完全相同的评测协议下，τ=0.10 是 238/274，也就是 86.86%。相对 Vanilla OPD 提升 7.30 个百分点，相对 TCOD-F2B 提升 2.19 个百分点。Seen 和 Unseen 分别是 87.86% 和 85.82%，说明提升没有只集中在训练分布内任务。

## 10 为什么有效与证据边界（70–90 秒）

机制上，随着 Student 变强，frontier 触发率从训练早期的 48% 降到后期 9%，而平均可蒸馏 horizon 从 20 增长到 28.2。也就是说，这个方法早期主动过滤，后期会自然接近 full-loss，而不是永远截断。我们目前最稳妥的结论是：相对熵漂移能够作为 trajectory-level 筛选信号，但关键是过滤强度和状态覆盖的平衡。边界也需要讲清楚：优势最强出现在最终 checkpoint，不是全程稳定领先；对 TCOD 的 2.19 个百分点在单种子配对检验下尚未显著。下一步会做多种子复测，并把 hard truncation 扩展到连续权重和 teacher compute allocation。
"""


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_clickable_box(slide, box_px: tuple[int, int, int, int], url: str) -> None:
    x1, y1, x2, y2 = box_px
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1 / 120),
        Inches(y1 / 120),
        Inches((x2 - x1) / 120),
        Inches((y2 - y1) / 120),
    )
    shape.fill.background()
    shape.line.fill.background()
    shape.click_action.hyperlink.address = url


def make_contact_sheet(paths: list[Path], out: Path) -> None:
    thumb_w, thumb_h = 640, 360
    sheet = Image.new("RGB", (thumb_w * 2 + 60, thumb_h * 5 + 120), "#E8E4E6")
    d = ImageDraw.Draw(sheet)
    draw_text(d, (30, 24), "面试自我介绍 PPT · 10 页预览", 28, PLUM, bold=True)
    for i, p in enumerate(paths):
        im = Image.open(p).resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        x = 20 + (i % 2) * (thumb_w + 20)
        y = 80 + (i // 2) * (thumb_h + 10)
        sheet.paste(im, (x, y))
    sheet.save(out, quality=94)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

    slide_paths: list[Path] = []
    for i, builder in enumerate(SLIDE_BUILDERS, start=1):
        im = builder()
        path = PREVIEW_DIR / f"slide_{i:02d}.png"
        im.save(path, optimize=True)
        slide_paths.append(path)

    prs = Presentation(str(REFERENCE))
    remove_all_slides(prs)
    blank = prs.slide_layouts[0]
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for idx, path in enumerate(slide_paths):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        for box, url in HYPERLINKS.get(idx, []):
            add_clickable_box(slide, box, url)

    out_pptx = OUT_DIR / "孙家兴_面试自我介绍_OPD项目.pptx"
    prs.save(out_pptx)
    (OUT_DIR / "孙家兴_面试自我介绍_讲稿.md").write_text(SPEAKER_NOTES, encoding="utf-8")
    make_contact_sheet(slide_paths, OUT_DIR / "preview_contact_sheet.png")
    print(out_pptx)


if __name__ == "__main__":
    main()
